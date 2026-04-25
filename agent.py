# COM6104 Topics in Data Science and AI - Agentic AI Project
# Full Implementation with UI, MCP Standard Tools, Short-term Memory, Multi-step Reasoning
import tkinter as tk
from tkinter import scrolledtext, END, messagebox
import ollama
from pptx import Presentation
import time

# Memory Module 
# Short-term memory system for recording interaction history
class AgentMemory:
    def __init__(self):
        self.history = []
        self.max_records = 10

    # Store user command and agent response
    def add_record(self, role, content):
        self.history.append({"role": role, "content": content})
        if len(self.history) > self.max_records:
            self.history.pop(0)

    # Get formatted memory history
    def get_memory_history(self):
        formatted_history = ""
        for record in self.history:
            formatted_history += f"{record['role']}: {record['content']}\n"
        return formatted_history

    # Clear all memory records
    def clear_memory(self):
        self.history = []

# MCP Standard Tool 1: Document Processing 
# Tool for reading and summarizing TXT and PPT files
class DocumentTool:
    def __init__(self):
        self.tool_name = "DocumentProcessor"
        self.tool_type = "Document_Analysis"
        self.tool_description = "Read text from TXT/PPTX files and generate content summary"

    # Read content from TXT file
    def read_txt_file(self, file_path="test.txt"):
        try:
            with open(file_path, "r", encoding="utf-8") as file:
                return file.read()
        except Exception as error:
            return f"Failed to read TXT file: {str(error)}"

    # Read text content from PPTX file
    def read_ppt_file(self, file_path="test.pptx"):
        try:
            presentation = Presentation(file_path)
            full_text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text += shape.text + " "
            return full_text
        except Exception as error:
            return f"Failed to read PPT file: {str(error)}"

    # Generate summary using local LLM
    def generate_summary(self, content):
        prompt = f"Summarize the following content in 2 to 3 short sentences. Do not ask questions. Content: {content}"
        response = ollama.chat(
            model="qwen2.5:0.5b",
            messages=[{"role": "user", "content": prompt}]
        )
        return response["message"]["content"]

# MCP Standard Tool 2: Task Management 
# Tool for task list management and countdown timer
class TaskManagerTool:
    def __init__(self):
        self.tool_name = "TaskManager"
        self.tool_type = "Task_Automation"
        self.tool_description = "Manage task list and execute countdown timer"
        self.task_list = []

    # Add new task to list
    def add_new_task(self, task_content):
        create_time = time.ctime()
        self.task_list.append({
            "task": task_content,
            "created_time": create_time
        })
        return f"Task added successfully: {task_content}"

    # Show all tasks in the list
    def show_all_tasks(self):
        if len(self.task_list) == 0:
            return "No tasks available in the list."
        task_output = "Task List:\n"
        for index, task_item in enumerate(self.task_list, 1):
            task_output += f"{index}. {task_item['task']} - Created: {task_item['created_time']}\n"
        return task_output

    # Countdown timer function
    def run_timer(self, seconds_input):
        try:
            total_seconds = int(seconds_input)
            time.sleep(total_seconds)
            return f"Timer completed. {total_seconds} seconds have passed."
        except ValueError:
            return "Invalid input. Please enter a valid number."

# Main agent for reasoning, tool selection and command execution
class AutonomousAgent:
    def __init__(self):
        self.memory = AgentMemory()
        self.document_tool = DocumentTool()
        self.task_tool = TaskManagerTool()

    # Parse command and execute corresponding functions
    def execute_command(self, user_command):
        self.memory.add_record("User", user_command)
        command = user_command.strip().lower()
        execution_result = ""

        # Document tool functions
        if command == "read txt":
            file_content = self.document_tool.read_txt_file()
            execution_result = f"TXT File Content:\n{file_content}"

        elif command == "summarize txt":
            file_content = self.document_tool.read_txt_file()
            summary = self.document_tool.generate_summary(file_content)
            execution_result = f"TXT File Summary:\n{summary}"

        elif command == "read ppt":
            file_content = self.document_tool.read_ppt_file()
            execution_result = f"PPT File Content:\n{file_content}"

        elif command == "summarize ppt":
            file_content = self.document_tool.read_ppt_file()
            summary = self.document_tool.generate_summary(file_content)
            execution_result = f"PPT File Summary:\n{summary}"

        # Task tool functions
        elif command.startswith("add task"):
            task_content = user_command.replace("add task", "").strip()
            execution_result = self.task_tool.add_new_task(task_content)

        elif command == "list tasks":
            execution_result = self.task_tool.show_all_tasks()

        elif command.startswith("timer"):
            seconds_value = user_command.replace("timer", "").strip()
            execution_result = self.task_tool.run_timer(seconds_value)

        # Memory functions
        elif command == "show memory":
            execution_result = f"Short-term Memory History:\n{self.memory.get_memory_history()}"

        elif command == "clear memory":
            self.memory.clear_memory()
            execution_result = "Short-term memory has been cleared."

        # Help information
        else:
            execution_result = ("Available Commands:\n"
                               "read txt - Read content from test.txt\n"
                               "summarize txt - Generate summary for test.txt\n"
                               "read ppt - Read content from test.pptx\n"
                               "summarize ppt - Generate summary for test.pptx\n"
                               "add task [task content] - Add a new task\n"
                               "list tasks - Show all tasks\n"
                               "timer [seconds] - Run countdown timer\n"
                               "show memory - Show interaction history\n"
                               "clear memory - Clear all memory records")

        self.memory.add_record("Agent", execution_result)
        return execution_result

class AgentUserInterface:
    def __init__(self, root_window):
        self.agent = AutonomousAgent()
        self.root = root_window
        self.root.title("COM6104 Autonomous AI Agent")
        self.root.geometry("850x650")
        self.build_interface()

    def build_interface(self):
        # Top control panel
        top_panel = tk.Frame(self.root)
        top_panel.pack(fill="x", padx=10, pady=5)
        tk.Label(top_panel, text="COM6104 Autonomous AI Agent", font=("Arial", 14, "bold")).pack(side="left")
        tk.Button(top_panel, text="Help", command=self.show_help).pack(side="right", padx=5)
        tk.Button(top_panel, text="Clear", command=self.clear_all_panel).pack(side="right", padx=5)

        # Command input area
        tk.Label(self.root, text="Command Input:", font=("Arial", 11)).pack(anchor="w", padx=10)
        self.command_input = scrolledtext.ScrolledText(self.root, height=2, width=100)
        self.command_input.pack(padx=10, pady=3)
        tk.Button(self.root, text="Run Command", bg="#4285F4", fg="white",
                  font=("Arial", 10, "bold"), command=self.run_agent_command).pack(pady=3)

        # Output display area
        tk.Label(self.root, text="Execution Output:", font=("Arial", 11)).pack(anchor="w", padx=10)
        self.output_display = scrolledtext.ScrolledText(self.root, width=100, height=28)
        self.output_display.pack(padx=10, pady=5)

    # Run command and show result
    def run_agent_command(self):
        user_input = self.command_input.get("1.0", END).strip()
        if not user_input:
            messagebox.showwarning("Warning", "Please enter a valid command.")
            return
        self.output_display.insert(END, f"Command: {user_input}\n")
        result = self.agent.execute_command(user_input)
        self.output_display.insert(END, f"Result: {result}\n\n")
        self.command_input.delete("1.0", END)

    # Clear all display content
    def clear_all_panel(self):
        self.output_display.delete("1.0", END)
        self.command_input.delete("1.0", END)

    # Show help information
    def show_help(self):
        help_info = ("COM6104 Autonomous AI Agent\n\n"
                     "This agent implements two MCP standard tools:\n"
                     "1. Document Processor: Read and summarize TXT/PPT files\n"
                     "2. Task Manager: Manage tasks and run countdown timer\n\n"
                     "Core Features: Short-term memory, multi-step reasoning, local execution\n"
                     "All codes are managed on GitHub.")
        messagebox.showinfo("Project Help", help_info)

if __name__ == "__main__":
    main_window = tk.Tk()
    app_interface = AgentUserInterface(main_window)
    main_window.mainloop()